import base64
import os
from datetime import datetime
from typing import List, Dict, Any, Union
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
import json
import re


def generate_id(path: str) -> str:
    return base64.urlsafe_b64encode(path.encode()).decode().rstrip("=")

def generate_title_from_content(content: str) -> str:
    # Assume the title is the first line of the content
    title = content.split('\n', 1)[0].strip()
    return title if title else ''

def read_pdf(file_path: str) -> str:
    content = ""
    with open(file_path, 'rb') as file:
        reader = PdfReader(file)
        for page in reader.pages:
            page_text = page.extract_text() or ""
            content += page_text + "\n"
    return content

def read_docx(file_path: str) -> str:
    doc = Document(file_path)
    content = ""
    for para in doc.paragraphs:
        content += para.text + "\n"
    return content

def read_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    content = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content += shape.text + "\n"
    return content

def read_html(file_path: str) -> str:
    with open(file_path, 'r') as file:
        content = file.read()
        soup = BeautifulSoup(content, "html.parser")
        text = soup.get_text()
    return text

def read_json(file_path: str) -> str:
    def format_json(obj: Any) -> str:
        """Recursively format JSON objects into a string."""
        if isinstance(obj, dict):
            return json.dumps({k: format_json(v) for k, v in obj.items()}, indent=2)
        elif isinstance(obj, list):
            return "\n".join(format_json(item) for item in obj)
        else:
            return json.dumps(obj)
    with open(file_path, 'r', encoding='utf-8') as file:
        data = json.load(file)
    
    return format_json(data)

def get_file_data(folder_path: str, base_path: str) -> List[Dict[str, Any]]:
    data_list = []
    for file_name in os.listdir(folder_path):
        file_path = os.path.join(folder_path, file_name)
        if os.path.isfile(file_path):
            if file_name.endswith('.pdf'):
                content = read_pdf(file_path)
            elif file_name.endswith('.docx'):
                content = read_docx(file_path)
            elif file_name.endswith('.pptx'):
                content = read_pptx(file_path)
            elif file_name.endswith('.html'):
                content = read_html(file_path)
            else:
                with open(file_path, 'r', encoding='utf-8') as file:
                    content = file.read().strip()
            
            relative_path = os.path.relpath(file_path, base_path)
            data_list.append({
                'id': generate_id(file_path),
                'content': content,
                'title': relative_path,
                'lastUpdated': datetime.fromtimestamp(os.path.getmtime(file_path)),
                'url': None
            })
        elif os.path.isdir(file_path):
            data_list.extend(get_file_data(file_path, base_path))  # Recursive call for directories
    return data_list

def fetch_url_content(url: str) -> str:
    """Fetch content from a URL and return it as a string."""
    try:
        response = requests.get(url)
        response.raise_for_status()  # Raise an error for HTTP issues
        content = response.text
        text = BeautifulSoup(content, "html.parser").get_text()        
        return text
    except requests.RequestException as e:
        print(f"Error fetching {url}: {e}")
        return ""

def add_urls_content(data_list: List[Dict[str, Any]], urls_file_path: str) -> None:
    """Read URLs from a file, fetch page content, and update the data_list."""
    if not os.path.isfile(urls_file_path):
        print(f"URLs file not found: {urls_file_path}")
        return

    with open(urls_file_path, 'r') as urls_file:
        urls = urls_file.readlines()

    for url in urls:
        url = url.strip()
        if url:
            content = fetch_url_content(url)
            data_list.append({
                'id': generate_id(url),
                'content': content,
                'title': generate_title_from_content(content),
                'lastUpdated': datetime.now(),
                'url': url
            })

def get_all_data() -> List[Dict[str, Any]]:
    base_folder_path = os.path.join(os.path.dirname(__file__), '../data')
    data_list = get_file_data(base_folder_path, base_folder_path)
    urls_file_path = os.path.join(os.path.dirname(__file__), 'URL.txt')
    add_urls_content(data_list, urls_file_path)
    return data_list